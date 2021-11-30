from pptx import Presentation
import os
import re
import sys
import json

pattern = """doi.org[a-zA-Z0-9\.\\\/\:\[\]]*"""
matcher = re.compile(pattern)

def read_ppt(file):
    dois = []
    prs = Presentation(file)
    for slide in prs.slides:
        isbib = False
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if '/bibliograph' in shape.text.lower():
                    isbib = True

        if isbib:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "doi.org" in(shape.text):
                        lines = shape.text.split('\n')
                        for l in lines:
                            d = matcher.findall(l)
                            if len(d)>0:
                                dois.append((d[0], file,l))
    return dois

    
def read_folder(loc = r"C:\Users\nikolay.tchakarov\Data\PPT", recursive=True):
    dois = []
    for file in os.listdir(loc):
        if ".pptx" in file:
            dois += read_ppt(os.path.join(loc,file))
        if os.path.isdir(os.path.join(loc,file)) and recursive:
            dois += read_folder(os.path.join(loc,file), recursive)
    return dois

if __name__=='__main__':
    loc = sys.argv[1]
    dois = read_folder(loc)

    home = os.environ.get('HOMEPATH')
    savedir = os.path.join(home,'dois')
    if not os.path.exists(savedir):
        os.mkdir(savedir)

    with open(os.path.join(savedir,'dois.json'),'w') as f:
        json.dump(dois,f,indent=3)
        f.close()

    with open(os.path.join(savedir,'dois.csv'),'w') as f:       
        [f.write(','.join(l[:2])+'\n') for l in dois]
        f.close()
